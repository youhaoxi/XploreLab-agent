import json
import logging

from ppt_template_model import TYPE_MAP, parse_json
from pptx import Presentation
from utils import delete_slide, delete_slide_range, duplicate_slide, move_slide


def fill_template(template_path, output_path, json_data):
    """
    Fill the template with the given json data.
    """
    prs = Presentation(template_path)
    data = json.loads(json_data)
    slides = parse_json(data)
    for slide in slides:
        if slide.type == "title":
            title_slide = prs.slides[TYPE_MAP[slide.type]]
            slide.render(title_slide)
        elif slide.type == "acknowledgement":
            acknowledgement_slide = prs.slides[TYPE_MAP[slide.type]]
            slide.render(acknowledgement_slide)
        else:
            new_slide = duplicate_slide(prs, prs.slides[TYPE_MAP[slide.type]])
            slide.render(new_slide)

    delete_slide_range(prs, range(2, 12))
    delete_slide(prs, 0)
    move_slide(prs, 1, len(prs.slides) - 1)
    prs.save(output_path)


def extract_json(content):
    """
    Extract the json data from the given content.
    """
    # extract content within "```json" and "```"
    json_data = content.split("```json")[1].split("```")[0]
    return json_data


if __name__ == "__main__":
    import argparse
    import datetime

    parser = argparse.ArgumentParser()
    parser.add_argument("-t", "--template", type=str, default="template/template_ori.pptx")
    default_output_filename = f"output-{datetime.datetime.now().strftime('%Y-%m-%d-%H-%M-%S')}.pptx"
    parser.add_argument("-o", "--output", type=str, default=default_output_filename)
    parser.add_argument("-i", "--input", type=str, required=True)
    parser.add_argument("--cache_dir", type=str, default=".temp")
    args = parser.parse_args()

    # set env var UTU_PPT_CACHE_DIR
    import os

    os.environ["UTU_PPT_CACHE_DIR"] = args.cache_dir

    logging.basicConfig(level=logging.INFO)
    template = args.template
    output = args.output
    input_json = args.input
    with open(input_json) as f:
        content = f.read()
    json_data = extract_json(content)
    fill_template(template, output, json_data)
