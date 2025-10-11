import copy
import logging
import random

import matplotlib
from pptx.enum.shapes import MSO_SHAPE_TYPE

# rgb colors from a color scheme
_color_palette = [
    (178, 34, 34),  # Brick Red
    (46, 139, 87),  # Sea Green
    (70, 130, 180),  # Steel Blue
    (210, 180, 140),  # Tan
    (147, 112, 219),  # Medium Purple
    (255, 165, 0),  # Orange (desaturated)
    (72, 209, 204),  # Medium Turquoise
    (205, 92, 92),  # Indian Red
    (106, 90, 205),  # Slate Blue
    (238, 130, 238),  # Violet
    (60, 179, 113),  # Medium Sea Green
    (100, 149, 237),  # Cornflower Blue
    (218, 165, 32),  # Goldenrod
    (199, 21, 133),  # Medium Violet Red
    (65, 105, 225),  # Royal Blue
]


def inspect_ppt(prs):
    """
    Inspect the given presentation.
    """
    for slide in prs.slides:
        inspect_slide(slide)


def inspect_slide(slide):
    """
    Inspect the given slide layout.
    """

    def _inspect_shape_list(shapes, indent=0):
        for shape in shapes:
            print(" " * indent + shape.name)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _inspect_shape_list(shape.shapes, indent + 2)

    _inspect_shape_list(slide.shapes)


def to_svg(slide, prs, svg_filename="test.svg"):
    """
    Convert the given slide layout to an SVG file.
    """

    def _to_svg_box(shapes, svg_box):
        for shape in shapes:
            svg_box.append(
                {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                    "shape_type": shape.shape_type,
                }
            )
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _to_svg_box(shape.shapes, svg_box)

    def render_svg(svg_box):
        width = prs.slide_width.inches
        height = prs.slide_height.inches
        fig, ax = matplotlib.pyplot.subplots(figsize=(width, height))
        ax.set_xlim(0, width)
        ax.set_ylim(height, 0)
        for box in svg_box:
            left, top, width, height = box["left"].inches, box["top"].inches, box["width"].inches, box["height"].inches
            picked_color = random.choice(_color_palette)
            color = (picked_color[0] / 255, picked_color[1] / 255, picked_color[2] / 255, 0.3)
            ax.add_patch(matplotlib.patches.Rectangle((left, top), width, height, color=color))
        ax.set_axis_off()
        matplotlib.pyplot.subplots_adjust(left=0, right=1, top=1, bottom=0)
        matplotlib.pyplot.margins(0)
        ax.set_xmargin(0)
        ax.set_ymargin(0)
        fig.savefig(svg_filename, bbox_inches="tight", pad_inches=0)

    svg_box = []
    _to_svg_box(slide.shapes, svg_box)
    render_svg(svg_box)

    return svg_box


def delete_shape(shape):
    """
    Delete the given shape.
    """
    parent = shape.element.getparent()
    parent.remove(shape.element)


def find_shape_with_name(shapes, name, depth=0):
    """
    Find the shape with the given name in the given shapes.
    """
    if depth == 0:
        logging.info(f"Finding shape with name: {name}")
    for shape in shapes:
        if shape.name == name:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found = find_shape_with_name(shape.shapes, name, depth + 1)
            if found:
                return found
    return None


def find_shape_with_name_except(shapes, name, depth=0):
    """
    Find the shape with the given name in the given shapes, except the shape with the given name.
    """
    if depth == 0:
        logging.info(f"Finding shape with name: {name}")
    for shape in shapes:
        if shape.name == name:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found = find_shape_with_name(shape.shapes, name, depth + 1)
            if found:
                return found
    raise Exception(f"Shape with name {name} not found")


def duplicate_slide(prs, slide):
    """
    Copy the given slide to the end of the presentation.

    Args:
        prs: Presentation object
        slide: Slide object to be copied
    Returns:
        New created Slide object
    """
    # Create new slide, use the same layout
    slide_layout = slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Copy all shapes
    for shape in slide.shapes:
        new_shape_element = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_shape_element, "p:extLst")

    return new_slide


def delete_slide_range(prs, index_range):
    """delete slides in the given index range
    Args:
        prs: Presentation object
        index_range: range of slide indices (0-based)
    """
    for index in reversed(index_range):
        delete_slide(prs, index)


def delete_slide(prs, index):
    """delete slide at the given index
    Args:
        prs: Presentation object
        index: slide index (0-based)

    Raises:
        IndexError: when index out of range
    """
    if index < 0 or index >= len(prs.slides):
        raise IndexError(f"Slide index {index} out of range (0-{len(prs.slides) - 1})")

    xml_slides = prs.slides._sldIdLst
    xml_slides.remove(xml_slides[index])


def move_slide(prs, old_index, new_index):
    """move slide from old_index to new_index"""
    xml_slides = prs.slides._sldIdLst
    # get the element to move
    slide_element = xml_slides[old_index]
    # remove the element from old position
    xml_slides.remove(slide_element)
    # insert to new position
    xml_slides.insert(new_index, slide_element)
    return prs


def replace_picture_keep_format(slide, old_picture_index, new_image_path):
    """
    Replace the picture at the given index with a new image while keeping all format attributes.
    """
    target_shape = slide.shapes[old_picture_index]

    # Save all properties
    properties = {
        "left": target_shape.left,
        "top": target_shape.top,
        "width": target_shape.width,
        "height": target_shape.height,
        "rotation": target_shape.rotation if hasattr(target_shape, "rotation") else 0,
    }

    # Save shadow, border, etc. effects (if exist)
    _shadow = None
    if hasattr(target_shape, "shadow"):
        _shadow = target_shape.shadow

    # Get the position in shapes collection
    _shape_index = None
    for i, shape in enumerate(slide.shapes):
        if shape == target_shape:
            _shape_index = i
            break

    # Remove the original picture
    sp = target_shape._element
    sp.getparent().remove(sp)

    # Insert the new picture
    new_picture = slide.shapes.add_picture(
        new_image_path, properties["left"], properties["top"], properties["width"], properties["height"]
    )

    # Restore rotation
    if properties["rotation"] != 0:
        new_picture.rotation = properties["rotation"]

    # Restore shadow, border, etc. effects
    if _shadow:
        new_picture.shadow = _shadow

    return new_picture
