from typing import Literal, Optional, Dict, Any, List
import json
from pydantic import BaseModel
from utils import delete_shape, find_shape_with_name_except, duplicate_slide, inspect_slide
from pptx import Presentation
import logging
import traceback
import requests
import uuid
from PIL import Image

TYPE_MAP = {
    "content": 0,
    "title": 1,
    "section_title": 2,
    "two_column_content": 3,
    "items_page_3": 4,
    "items_page_4": 5,
    "items_page_4_alt": 6,
    
    "item_page_4_growth": 7,
    "items_page_4_key_point": 8,
    
    "item_page_6_sequence": 9,
    "items_page_4_img": 10,
    "item_page_5_sequence": 11,
    "acknowledgement": 12
}

class BaseContent(BaseModel):
    content_type: Literal["text", "image", "table"]

class Paragraph(BaseModel):
    text: str
    bullet: bool = False 
    level: int = 0

class Item(BaseModel):
    title: str # be very concise within 3 words, or 4 characters
    content: str # be very concise within 10 words

class TextContent(BaseContent):
    content_type: Literal["text"] = "text"
    paragraph: list[Paragraph] | str

class ImageContent(BaseContent):
    content_type: Literal["image"] = "image"
    image_url: str # absolute url
    caption: Optional[str] = None # be very concise within 20 words
    
class TableContent(BaseContent):
    content_type: Literal["table"] = "table"
    header: list[str]
    rows: list[list[str]]
    caption: Optional[str] = None # be very concise within 20 words
    n_rows: int # no more than 7
    n_cols: int # no more than 10
    
class Slide(BaseModel):
    type: Literal["title", "section_title", "content", "two_column_content", "items_page_3", "items_page_4", "items_page_4_key_point", "items_page_4_img", "item_page_4_growth", "item_page_6_sequence", "item_page_5_sequence", "acknowledgement"]
    
class TitlePage(Slide):
    """
    Title page
    
    note: should be the first page
    """
    type: Literal["title"] = "title"
    title: str = "" # be very concise within 7 words
    subtitle: str = "" # be very concise within 10 words
    author: str = ""
    date: str = ""
    ppt_type: str = "" # describe the type of the ppt, e.g., "work report"
    
    def render(self, slide):
        logging.info(f"===Rendering title page: {self.title}===")
        title_shape = find_shape_with_name_except(slide.shapes, "title")
        subtitle_shape = find_shape_with_name_except(slide.shapes, "subtitle")
        author_shape = find_shape_with_name_except(slide.shapes, "author")
        date_shape = find_shape_with_name_except(slide.shapes, "date")
        ppt_type_shape = find_shape_with_name_except(slide.shapes, "ppt_type")
        
        logging.info(f"title: {self.title}")
        handle_pure_text(self.title, title_shape, slide)
        logging.info(f"subtitle: {self.subtitle}")
        handle_pure_text(self.subtitle, subtitle_shape, slide)
        logging.info(f"author: {self.author}")
        handle_pure_text(self.author, author_shape, slide)
        logging.info(f"date: {self.date}")
        handle_pure_text(self.date, date_shape, slide)
        logging.info(f"ppt_type: {self.ppt_type}")
        handle_pure_text(self.ppt_type, ppt_type_shape, slide)
    
class SectionTitlePage(Slide):
    """
    Section Title Page
    
    note: should not follow the title page immediately
    """
    type: Literal["section_title"] = "section_title"
    title: str = "" # be very concise within 7 words
    subtitle: str = "" # be very concise within 10 words
    section_no: int # section number
    
    def render(self, slide):
        logging.info(f"===Rendering section title page: {self.title}===")
        title_shape = find_shape_with_name_except(slide.shapes, "title")
        subtitle_shape = find_shape_with_name_except(slide.shapes, "subtitle")
        section_no_shape = find_shape_with_name_except(slide.shapes, "section_no")
        
        logging.info(f"title: {self.title}")
        handle_pure_text(self.title, title_shape, slide)
        logging.info(f"subtitle: {self.subtitle}")
        handle_pure_text(self.subtitle, subtitle_shape, slide)
        handle_pure_text(f"{self.section_no:02d}", section_no_shape, slide)

class ContentPage(Slide):
    """
    Single column content page, with the most basic layout.
    """
    type: Literal["content"] = "content"
    title: str = "" # be very concise within 7 words
    content: BaseContent
    
    def render(self, slide):
        logging.info(f"===Rendering content page: {self.title}===")
        title_shape = find_shape_with_name_except(slide.shapes, "title")
        content_shape = find_shape_with_name_except(slide.shapes, "content")
        
        logging.info(f"title: {self.title}")
        handle_pure_text(self.title, title_shape, slide)
        logging.info(f"content: {self.content}")
        handle_content(self.content, content_shape, slide)
    
class TwoColumnContentPage(Slide):
    """
    Two column content pages
    """
    type: Literal["two_column_content"] = "two_column_content"
    title: str = "" # be very concise within 7 words
    content: list[BaseContent]
    
    def render(self, slide):
        logging.info(f"===Rendering two column content page: {self.title}===")
        title_shape = find_shape_with_name_except(slide.shapes, "title")
        content1_shape = find_shape_with_name_except(slide.shapes, "content1")
        content2_shape = find_shape_with_name_except(slide.shapes, "content2")
        
        logging.info(f"title: {self.title}")
        handle_pure_text(self.title, title_shape, slide)
        logging.info(f"content1: {self.content[0]}")
        handle_content(self.content[0], content1_shape, slide)
        logging.info(f"content2: {self.content[1]}")
        handle_content(self.content[1], content2_shape, slide)

class ItemsPage3(Slide):
    """
    Page with three logically parallel items
    """
    type: Literal["items_page_3"] = "items_page_3"
    title: str = "" # be very concise within 7 words
    items: list[Item] # exactly 3 items
    
    def render(self, slide):
        logging.info(f"===Rendering items page 3: {self.title}===")
        title_shape = find_shape_with_name_except(slide.shapes, "title")
        
        logging.info(f"title: {self.title}")
        handle_pure_text(self.title, title_shape, slide)
        
        for ind, item in enumerate(self.items):
            logging.info(f"item {ind}: {item}")
            handle_item(item, ind, slide)

class ItemsPage4(Slide):
    """
    Page with four logically parallel items
    """
    type: Literal["items_page_4"] = "items_page_4"
    title: str = "" # be very concise within 7 words
    items: list[Item] # exactly 4 items
    
    def render(self, slide):
        logging.info(f"===Rendering items page 4: {self.title}===")
        title_shape = find_shape_with_name_except(slide.shapes, "title")
        
        logging.info(f"title: {self.title}")
        handle_pure_text(self.title, title_shape, slide)
        
        for ind, item in enumerate(self.items):
            logging.info(f"item {ind}: {item}")
            handle_item(item, ind, slide)

class ItemsPage4KeyPoint(Slide):
    """
    This slide is designed to present four key elements, stages, or pillars in a clean, balanced layout. Each element is represented by a customizable label — which can be:
    A letter (e.g., S-W-O-T)
    A Chinese character (e.g., 增 / 删 / 改 / 查)
    """
    type: Literal["items_page_4_key_point"] = "items_page_4_key_point"
    title: str = "" # be very concise within 7 words
    label: Optional[list[str]] = None # be exactly one character/letter/digit
    items: list[Item] # exactly 4 items
    
    def render(self, slide):
        logging.info(f"===Rendering items page 4 key point: {self.title}===")
        title_shape = find_shape_with_name_except(slide.shapes, "title")
        
        logging.info(f"title: {self.title}")
        handle_pure_text(self.title, title_shape, slide)
        
        for ind, item in enumerate(self.items):
            logging.info(f"item {ind}: {item}")
            handle_item(item, ind, slide)

class ItemsPage4Img(Slide):
    """
    A process or framework centered around data/analysis, with 4 connected stages represented by icons — useful for explaining systems, strategies, or operations.
    """
    type: Literal["items_page_4_img"] = "items_page_4_img"
    title: str = "" # be very concise within 7 words
    items: list[Item] # exactly 4 items
    
    def render(self, slide):
        logging.info(f"===Rendering items page 4 img: {self.title}===")
        title_shape = find_shape_with_name_except(slide.shapes, "title")
        
        logging.info(f"title: {self.title}")
        handle_pure_text(self.title, title_shape, slide)
        
        for ind, item in enumerate(self.items):
            logging.info(f"item {ind}: {item}")
            handle_item(item, ind, slide)

class ItemPage4Growth(Slide):
    """
    This slide visually represents growth over time, using an upward-sloping bar chart. Four key stages or drivers of growth are highlighted on the left, each with an icon (home, battery, lock, chat bubbles) and placeholder text. The design suggests that these factors contribute progressively to overall growth, making it ideal for showing business expansion, performance metrics, or KPIs.
    """
    type: Literal["item_page_4_growth"] = "item_page_4_growth"
    title: str = "" # be very concise within 7 words
    items: list[Item]  # Exactly 4 items. The first item (index 0) corresponds to the topmost / highest point on the growth curve; subsequent items map downward in order along the ascending bar chart.
    
    def render(self, slide):
        logging.info(f"===Rendering item page 4 growth: {self.title}===")
        title_shape = find_shape_with_name_except(slide.shapes, "title")
        
        logging.info(f"title: {self.title}")
        handle_pure_text(self.title, title_shape, slide)
        
        for ind, item in enumerate(self.items):
            logging.info(f"item {ind}: {item}")
            handle_item(item, ind, slide)
            
class ItemPage6Sequence(Slide):
    """
    This slide illustrates a step-by-step process or timeline. Six stages are connected by a dotted line, moving from left to right. Each stage has an icon (database, fork/merge, gears, lightbulb, pen, etc.) and placeholder text. It’s perfect for explaining workflows, project phases, or a logical progression of ideas or actions.
    """
    type: Literal["item_page_6_sequence"] = "item_page_6_sequence"
    title: str = "" # be very concise within 7 words
    items: list[Item] # exactly 6 items
    
    def render(self, slide):
        logging.info(f"===Rendering item page 6 sequence: {self.title}===")
        title_shape = find_shape_with_name_except(slide.shapes, "title")
        
        logging.info(f"title: {self.title}")
        handle_pure_text(self.title, title_shape, slide)
        
        for ind, item in enumerate(self.items):
            logging.info(f"item {ind}: {item}")
            handle_item(item, ind, slide)

class ItemPage5Sequence(Slide):
    """
    This slide illustrates a step-by-step process or timeline. Six stages are connected by a dotted line, moving from left to right. Each stage has an icon (database, fork/merge, gears, lightbulb, pen, etc.) and placeholder text. It’s perfect for explaining workflows, project phases, or a logical progression of ideas or actions.
    """
    type: Literal["item_page_5_sequence"] = "item_page_5_sequence"
    title: str = "" # be very concise within 7 words
    items: list[Item] # exactly 5 items
    
    def render(self, slide):
        logging.info(f"===Rendering item page 5 sequence: {self.title}===")
        title_shape = find_shape_with_name_except(slide.shapes, "title")
        
        logging.info(f"title: {self.title}")
        handle_pure_text(self.title, title_shape, slide)
        
        for ind, item in enumerate(self.items):
            logging.info(f"item {ind}: {item}")
            handle_item(item, ind, slide)

class AcknowledgementPage(Slide):
    """
    Acknowledgement page
    
    note: should be the last page
    """
    type: Literal["acknowledgement"] = "acknowledgement"
    title: str = "" # just place "Thank you"/"谢谢"/"END"/"FIN" here
    subtitle: str = "" # be very concise within 10 words
    author: str = ""
    date: str = ""
    ppt_type: str = "" # describe the type of the ppt, e.g., "work report"
    
    def render(self, slide):
        logging.info(f"===Rendering acknowledgement page: {self.title}===")
        title_shape = find_shape_with_name_except(slide.shapes, "title")
        subtitle_shape = find_shape_with_name_except(slide.shapes, "subtitle")
        author_shape = find_shape_with_name_except(slide.shapes, "author")
        date_shape = find_shape_with_name_except(slide.shapes, "date")
        ppt_type_shape = find_shape_with_name_except(slide.shapes, "ppt_type")
        
        logging.info(f"title: {self.title}")
        handle_pure_text(self.title, title_shape, slide)
        logging.info(f"subtitle: {self.subtitle}")
        handle_pure_text(self.subtitle, subtitle_shape, slide)
        logging.info(f"author: {self.author}")
        handle_pure_text(self.author, author_shape, slide)
        logging.info(f"date: {self.date}")
        handle_pure_text(self.date, date_shape, slide)
        logging.info(f"ppt_type: {self.ppt_type}")
        handle_pure_text(self.ppt_type, ppt_type_shape, slide)

def parse_json(json_data: Dict[str, Any]) -> List[Slide]:
    slides = []
    for slide_data in json_data['slides']:
        slide_type = slide_data['type']
        
        if slide_type == 'title':
            slides.append(TitlePage(**slide_data))
        elif slide_type == 'section_title':
            slides.append(SectionTitlePage(**slide_data))
        elif slide_type == 'content':
            content_page = ContentPage(**slide_data)
            content_data = slide_data['content']
            if content_data['content_type'] == 'text':
                content_page.content = TextContent(**content_data)
            elif content_data['content_type'] == 'image':
                content_page.content = ImageContent(**content_data)
            elif content_data['content_type'] == 'table':
                content_page.content = TableContent(**content_data)
            else:
                raise ValueError(f"Unsupported content type: {content_data['content_type']}")
            slides.append(content_page)
        elif slide_type == 'two_column_content':
            content_page = TwoColumnContentPage(**slide_data)
            content_data = slide_data['content']
            for i in range(2):
                if content_data[i]['content_type'] == 'text':
                    content_page.content[i] = TextContent(**content_data[i])
                elif content_data[i]['content_type'] == 'image':
                    content_page.content[i] = ImageContent(**content_data[i])
                elif content_data[i]['content_type'] == 'table':
                    content_page.content[i] = TableContent(**content_data[i])
                else:
                    raise ValueError(f"Unsupported content type: {content_data['content_type']}")
            slides.append(content_page)
        elif slide_type == 'items_page_3':
            slides.append(ItemsPage3(**slide_data))
        elif slide_type == 'items_page_4':
            slides.append(ItemsPage4(**slide_data))
        elif slide_type == 'items_page_4_key_point':
            slides.append(ItemsPage4KeyPoint(**slide_data))
        elif slide_type == 'items_page_4_img':
            slides.append(ItemsPage4Img(**slide_data))
        elif slide_type == 'item_page_4_growth' or slide_type == 'items_page_4_growth':
            slide_data['type'] = 'item_page_4_growth'
            slides.append(ItemPage4Growth(**slide_data))
        elif slide_type == 'item_page_6_sequence' or slide_type == 'items_page_6_sequence':
            slide_data['type'] = 'item_page_6_sequence'
            slides.append(ItemPage6Sequence(**slide_data))
        elif slide_type == 'acknowledgement':
            slides.append(AcknowledgementPage(**slide_data))
        elif slide_type == 'item_page_5_sequence' or slide_type == 'items_page_5_sequence':
            slide_data['type'] = 'item_page_5_sequence'
            slides.append(ItemPage5Sequence(**slide_data))
        else:
            raise ValueError(f"Unsupported slide type: {slide_type}")
    return slides

def download_image(url, base_dir="."):
    if url.startswith("http"):
        # download the image
        headers = {
            'Accept': 'image/*, */*'
        }
        response = requests.get(url, headers=headers)
        extension_name = url.split(".")[-1]
        if extension_name not in ["png", "jpg", "jpeg", "gif", "bmp", "webp"]:
            extension_name = "png"
        if response.status_code == 200:
            file_name = f"{base_dir}/{uuid.uuid4()}.{extension_name}"
            with open(file_name, "wb") as f:
                f.write(response.content)
            # get width and height
            image = Image.open(file_name)
            width, height = image.size
            return file_name, width, height
        else:
            raise Exception(f"Failed to download image: {url} {response.status_code}")
    raise Exception(f"Failed to download image: {url}")

def handle_pure_text(text: str, target_shape, slide):
    try:
        text_frame = target_shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                original_font = run.font
                run.text = text
    except Exception as e:
        logging.error(f"Failed to set text: {text} {e}")
        traceback.print_exc()
        raise Exception(f"Failed to set text: {text}")

def handle_image(image_url: str, target_shape, slide):
    left, top, width, height = target_shape.left, target_shape.top, target_shape.width, target_shape.height
    try:
        image_url, image_width, image_height = download_image(image_url)
    except Exception as e:
        logging.warning(f"Failed to download image: {image_url} {e}")
        traceback.print_exc()
        return
    
    # scale the image to fit the placeholder
    scale = min(width / image_width, height / image_height)
    image_width *= scale
    image_height *= scale
    # center the image
    left += (width - image_width) / 2
    top += (height - image_height) / 2
    slide.shapes.add_picture(image_url, left, top, image_width, image_height)
    
    # remove the placeholder
    delete_shape(target_shape)

def handle_table(table_content: TableContent, target_shape, slide):
    left, top, width, height = target_shape.left, target_shape.top, target_shape.width, target_shape.height
    table = slide.shapes.add_table(table_content.n_rows + 1, table_content.n_cols, left, top, width, height).table
    for i, header_text in enumerate(table_content.header):
        table.cell(0, i).text = header_text
    for i, row in enumerate(table_content.rows):
        for j, cell_text in enumerate(row):
            table.cell(i+1, j).text = cell_text
    table.auto_fit = True
    
    # remove the placeholder
    delete_shape(target_shape)

def handle_text_content(text_content: TextContent, target_shape, slide):
    if isinstance(text_content.paragraph, list):
        original_font = None
        for paragraph in target_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                original_font = run.font
                break
        
        text_frame = target_shape.text_frame
        
        text_frame.clear()
        
        for paragraph in text_content.paragraph:
            para = text_frame.add_paragraph()
            para.bullet = paragraph.bullet
            para.level = paragraph.level
            run = para.add_run()
            run.text = paragraph.text
            if original_font:
                run.font.name = original_font.name
                run.font.size = original_font.size
                run.font.bold = original_font.bold
                run.font.italic = original_font.italic
                if original_font.color.type == 1:
                    run.font.color.rgb = original_font.color.rgb
    else:
        handle_pure_text(text_content.paragraph, target_shape, slide)
    
def handle_content(content: BaseContent, target_shape, slide):
    if content.content_type == "text":
        handle_text_content(content, target_shape, slide)
    elif content.content_type == "image":
        handle_image(content.image_url, target_shape, slide)
    elif content.content_type == "table":
        handle_table(content, target_shape, slide)

def handle_item(item: Item, item_index: int, slide, index_start_from_one=True):
    if index_start_from_one:
        item_index += 1
    item_title_name = f"item_title{item_index}"
    item_content_name = f"item_content{item_index}"
    
    item_title_shape = find_shape_with_name_except(slide.shapes, item_title_name)
    item_content_shape = find_shape_with_name_except(slide.shapes, item_content_name)
    
    handle_pure_text(item.title, item_title_shape, slide)
    handle_pure_text(item.content, item_content_shape, slide)

if __name__ == '__main__':
    # logging level
    logging.basicConfig(level=logging.INFO)
    
    # simple test
    with open("template_example.json", "r") as f:
        data = json.load(f)
    slides = parse_json(data)
    
    ppt = Presentation("template/template_ori.pptx")
    for slide in slides:
        if slide.type == "title":
            title_slide = ppt.slides[TYPE_MAP[slide.type]]
            slide.render(title_slide)
        elif slide.type == "acknowledgement":
            acknowledgement_slide = ppt.slides[TYPE_MAP[slide.type]]
            slide.render(acknowledgement_slide)
        else:
            new_slide = duplicate_slide(ppt, ppt.slides[TYPE_MAP[slide.type]])
            slide.render(new_slide)
    ppt.save("gen.pptx")
    